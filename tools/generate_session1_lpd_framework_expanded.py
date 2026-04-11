from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT_DIR = Path('deliverables')
OUT_DIR.mkdir(exist_ok=True)
OUT_PPTX = OUT_DIR / 'Session_1_LPD_Framework_Expanded.pptx'
OUT_NOTES = OUT_DIR / 'Session_1_LPD_Framework_Expanded_Speaker_Notes.md'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(244, 245, 241)
DARK = RGBColor(26, 36, 32)
GREEN = RGBColor(79, 101, 74)
GREEN_2 = RGBColor(132, 150, 121)
LIGHT = RGBColor(255, 255, 255)
TAN = RGBColor(233, 226, 210)
ACCENT = RGBColor(123, 96, 63)
GRAY = RGBColor(90, 96, 92)

TITLE_FONT = 'Arial'
BODY_FONT = 'Arial'


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def top(slide, title, subtitle=''):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = DARK
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.16), Inches(11.8), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = TITLE_FONT
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.color.rgb = LIGHT
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = BODY_FONT
        r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(214, 220, 213)


def footer(slide, text='Session 1 | Foundation of Leadership | Expanded from original LPD framework'):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(7.0), Inches(12.35), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_2
    line.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.02), Inches(8), Inches(0.2))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY


def textbox(slide, left, top_, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top_, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def bullets(slide, left, top_, width, height, items, size=18, color=DARK):
    tb = slide.shapes.add_textbox(left, top_, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = text
        r.font.name = BODY_FONT
        r.font.size = Pt(size - level * 2)
        r.font.color.rgb = color
    return tb


def card(slide, left, top_, width, height, header, items, header_color=GREEN, fill_color=LIGHT, size=16):
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top_, width, height)
    body.fill.solid()
    body.fill.fore_color.rgb = fill_color
    body.line.color.rgb = GREEN_2
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top_, width, Inches(0.42))
    head.fill.solid()
    head.fill.fore_color.rgb = header_color
    head.line.fill.background()
    textbox(slide, left + Inches(0.18), top_ + Inches(0.06), width - Inches(0.3), Inches(0.24), header, size=14, color=LIGHT, bold=True)
    bullets(slide, left + Inches(0.18), top_ + Inches(0.52), width - Inches(0.3), height - Inches(0.62), items, size=size)


def quote(slide, text, source):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.68))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TAN
    shape.line.color.rgb = RGBColor(213, 201, 170)
    textbox(slide, Inches(1.0), Inches(6.07), Inches(11.3), Inches(0.2), f'"{text}"', size=14, color=DARK, bold=True)
    textbox(slide, Inches(1.0), Inches(6.31), Inches(5), Inches(0.18), source, size=10, color=ACCENT)


def title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.85), Inches(0.18), Inches(5.0))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = GREEN_2; stripe.line.fill.background()
    textbox(slide, Inches(1.05), Inches(1.0), Inches(10.6), Inches(1.15), 'Session 1: The Foundation of Leadership', size=28, color=LIGHT, bold=True)
    textbox(slide, Inches(1.05), Inches(2.1), Inches(8.3), Inches(0.45), 'Influence and Trust | Expanded from the original LPD framework', size=18, color=RGBColor(220, 227, 219))
    textbox(slide, Inches(1.05), Inches(3.2), Inches(6.5), Inches(1.6), 'This version keeps the framework intact and builds it out with facilitator language, doctrine tie-ins, discussion prompts, and delivery notes.', size=18, color=RGBColor(231, 236, 230))
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(4.35), Inches(4.25), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(44, 57, 50); box.line.color.rgb = GREEN
    textbox(slide, Inches(8.45), Inches(4.68), Inches(3.7), Inches(0.6), 'Leadership in the Army is influence backed by trust.', size=18, color=LIGHT, bold=True)
    textbox(slide, Inches(8.45), Inches(5.38), Inches(3.0), Inches(0.2), 'Ready for final unit-level edits', size=11, color=RGBColor(210, 218, 208))


def purpose_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, 'How to Use This Expanded Framework', 'Built to preserve your original structure while giving you more usable content in the room.')
    card(slide, Inches(0.7), Inches(1.45), Inches(5.9), Inches(4.95), 'What this version does', [
        'Keeps the original Session 1 theme: Leadership foundations through influence and trust.',
        'Preserves the three anchor concepts from the spreadsheet: Lid, Influence, Solid Ground.',
        'Expands doctrine tie-ins, talk tracks, and discussion prompts for each concept.',
        'Adds facilitator-ready exercise guidance, timing, and closing prompts.'
    ], header_color=GREEN, fill_color=LIGHT, size=17)
    card(slide, Inches(6.8), Inches(1.45), Inches(5.8), Inches(4.95), 'Recommended use', [
        'Use the slides as your visual backbone, not a word-for-word script.',
        'Trim or combine slides depending on available time.',
        'Insert one or two examples from your formation before presenting.',
        'If presenting to junior Soldiers, simplify doctrine language and keep examples practical.'
    ], header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=17)
    textbox(slide, Inches(0.85), Inches(6.5), Inches(11.5), Inches(0.2), 'Best fit: a 45 to 60 minute LPD with discussion, one activity, and a clear closeout challenge.', size=12, color=GRAY, bold=True)
    footer(slide)


def overview_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, 'Original Framework, Expanded', 'This slide mirrors the source structure and shows how the brief is organized.')
    table = slide.shapes.add_table(4, 4, Inches(0.55), Inches(1.5), Inches(12.25), Inches(4.5)).table
    widths = [2.1, 2.55, 2.85, 4.75]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    headers = ['Session', 'Maxwell principle', 'Army reference', 'Expanded briefing focus']
    for i, text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = text
        cell.fill.solid(); cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = TITLE_FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = LIGHT
    rows = [
        ['Foundation of Leadership', 'Law of the Lid', 'ADP 6-22, FM 6-22', 'Why leader development raises unit effectiveness and why stagnation hurts the team.'],
        ['Foundation of Leadership', 'Law of Influence', 'ADP 6-22, FM 1', 'How leaders earn followership beyond title, authority, or positional power.'],
        ['Foundation of Leadership', 'Law of Solid Ground', 'AR 600-100, DA PAM 165-19', 'Why trust is the foundation for climate, standards, and team cohesion.']
    ]
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else RGBColor(249, 249, 245)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = BODY_FONT; r.font.size = Pt(12); r.font.color.rgb = DARK
    card(slide, Inches(0.85), Inches(6.15), Inches(11.35), Inches(0.68), 'Facilitator note', ['Present these ideas as reinforcing Army doctrine, not replacing it. Maxwell gives you language. Doctrine gives you legitimacy and standards.'], header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=14)
    footer(slide)


def concept_slide(title_txt, subtitle, left_header, left_items, right_header, right_items, quote_text, quote_src):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, title_txt, subtitle)
    card(slide, Inches(0.7), Inches(1.35), Inches(5.95), Inches(4.95), left_header, left_items, header_color=GREEN, fill_color=LIGHT, size=17)
    card(slide, Inches(6.85), Inches(1.35), Inches(5.75), Inches(4.95), right_header, right_items, header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=16)
    quote(slide, quote_text, quote_src)
    footer(slide)


def discussion_slide(title_txt, subtitle, framing, prompts, takeaways):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, title_txt, subtitle)
    card(slide, Inches(0.7), Inches(1.45), Inches(5.8), Inches(4.9), 'Facilitator framing', framing, header_color=GREEN, fill_color=LIGHT, size=17)
    card(slide, Inches(6.7), Inches(1.45), Inches(5.95), Inches(3.15), 'Discussion prompts', prompts, header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=16)
    card(slide, Inches(6.7), Inches(4.85), Inches(5.95), Inches(1.5), 'What you want them to leave with', takeaways, header_color=GREEN_2, fill_color=LIGHT, size=15)
    footer(slide)


def exercise_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, 'Exercise Expansion', 'These keep the activities from the spreadsheet but give you a more usable facilitation lane.')
    card(slide, Inches(0.7), Inches(1.4), Inches(5.9), Inches(5.15), 'Exercise 1 | The Lid self-assessment', [
        'Purpose: drive honest self-assessment and produce one development action.',
        'Method: have leaders score themselves against ADP 6-22 competencies or attributes.',
        'Ask three questions: What is my biggest lid? What resource can help? What will I do in 30 days?',
        'Use small-group sharing for goals, not personal weaknesses, to keep the tone professional.',
        'Close by asking for one volunteer takeaway and one concrete next step.'
    ], header_color=GREEN, fill_color=LIGHT, size=16)
    card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.15), 'Exercise 2 | Influence beyond rank role-play', [
        'Purpose: show the difference between authority and influence.',
        'Scenario: last-minute clean-up detail cancels early release.',
        'Leader task: gain buy-in without using "because I said so."',
        'Observer task: note what actually created willingness to follow.',
        'AAR questions: What language built buy-in? What shut it down? What would work in our unit?' 
    ], header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=16)
    footer(slide)


def icebreakers_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, 'Icebreakers from the Framework, Now with Recommended Use', 'Use only one. The goal is to loosen the room and transition into leadership discussion.')
    card(slide, Inches(0.65), Inches(1.35), Inches(3.0), Inches(4.95), 'Fastest', [
        'One Word for Leadership',
        '3 to 5 minutes',
        'Best for quick engagement and easy transition to trust and influence.'
    ], header_color=GREEN, fill_color=LIGHT, size=15)
    card(slide, Inches(3.95), Inches(1.35), Inches(3.0), Inches(4.95), 'Most flexible', [
        'Leadership If-Then Statement',
        '3 to 5 minutes',
        'Good when you want concise participation from everyone.'
    ], header_color=GREEN_2, fill_color=RGBColor(251, 250, 247), size=15)
    card(slide, Inches(7.25), Inches(1.35), Inches(2.75), Inches(4.95), 'Most fun', [
        'Two Truths and a Lie',
        '5 to 7 minutes',
        'Good for breaking tension in a new or mixed group.'
    ], header_color=ACCENT, fill_color=LIGHT, size=15)
    card(slide, Inches(10.25), Inches(1.35), Inches(2.4), Inches(4.95), 'Most reflective', [
        'Leadership Superpower or Most Valuable Lesson Learned',
        '7 to 15 minutes',
        'Best for experienced groups or deeper discussion.'
    ], header_color=DARK, fill_color=RGBColor(251, 250, 247), size=14)
    textbox(slide, Inches(0.8), Inches(6.45), Inches(11.6), Inches(0.24), 'Recommendation if you need a clean start: use One Word for Leadership, write responses on a board, then connect them to trust, development, and influence.', size=12, color=GRAY, bold=True)
    footer(slide)


def close_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, 'Suggested Closeout', 'This ending keeps the session actionable and brings the lesson back to the unit.')
    card(slide, Inches(0.7), Inches(1.45), Inches(5.85), Inches(4.9), 'Three closeout points', [
        'Your current leadership ability affects more people than just you.',
        'Influence is what turns compliance into commitment.',
        'Trust is built slowly and lost quickly, so leaders must guard it daily.'
    ], header_color=GREEN, fill_color=LIGHT, size=17)
    card(slide, Inches(6.8), Inches(1.45), Inches(5.8), Inches(4.9), 'Final challenge to the room', [
        'Identify one leadership lid you will raise in the next 30 days.',
        'Do one thing this week that builds influence without leaning on rank.',
        'Do one thing this week that visibly strengthens trust with your Soldiers.',
        'If possible, revisit commitments at the next LPD, counseling, or huddle.'
    ], header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=16)
    textbox(slide, Inches(0.85), Inches(6.45), Inches(11.0), Inches(0.22), 'Suggested final question: What will your Soldiers notice this week that proves this lesson changed something in you?', size=12, color=GRAY, bold=True)
    footer(slide)


def refs_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    top(slide, 'References and Final Edit Checklist', 'Sources from the spreadsheet plus the edits you should make before presenting')
    card(slide, Inches(0.7), Inches(1.45), Inches(5.8), Inches(4.95), 'References', [
        'John C. Maxwell, Leadership 101',
        'ADP 6-22, Army Leadership and the Profession',
        'FM 6-22, Developing Leaders',
        'FM 1, The Army',
        'AR 600-100, Army Profession and Leadership Policy',
        'DA PAM 165-19, Moral Leadership',
        'Supporting references listed in the source: ADP 1, DA PAM 600-25, TC 7-22.7, ADP 7-0'
    ], header_color=GREEN, fill_color=LIGHT, size=16)
    card(slide, Inches(6.8), Inches(1.45), Inches(5.8), Inches(4.95), 'Before you brief', [
        'Add your rank, unit, and date on the title slide.',
        'Swap in one example from your unit for each main concept.',
        'Adjust language for the maturity and rank of the audience.',
        'If short on time, cut one discussion slide and keep one exercise.',
        'If you want a stronger Army feel, add a unit crest or patch image in Keynote.'
    ], header_color=ACCENT, fill_color=RGBColor(251, 250, 247), size=16)
    footer(slide)


title_slide()
purpose_slide()
overview_slide()

concept_slide(
    'Law of the Lid | Expanded',
    'Original framework idea: Leadership ability determines a person\'s level of effectiveness.',
    'Doctrine tie-in and explanation', [
        'ADP 6-22 and FM 6-22 emphasize continuous development, reflection, and self-improvement.',
        'A leader\'s strengths or weaknesses do not stay personal. They shape the capability of the whole formation.',
        'When a leader stops learning, the team often feels it in standards, communication, and initiative.',
        'The practical Army question is not whether we have a lid. It is whether we are raising it on purpose.'
    ],
    'What to emphasize in the room', [
        'Development is not just PME. It includes counseling, feedback, repetition, reading, coaching, and humility.',
        'Leaders should identify one specific limitation, not a vague desire to get better.',
        'Command climate matters. If growth is punished, leaders stop trying.'
    ],
    'The lower an individual\'s ability to lead, the lower the lid on his potential. The higher the leadership, the greater the effectiveness.',
    'John C. Maxwell'
)

discussion_slide(
    'Law of the Lid | Facilitated Discussion',
    'Keep the discussion practical and grounded in what leaders see every week.',
    [
        'Frame it this way: every leader has a ceiling somewhere, and it usually shows up under stress.',
        'Ask for examples of common lids in junior leaders, team leaders, or squad leaders.',
        'Push toward action by asking how we deliberately raise that ceiling.'
    ],
    [
        'How does the Army\'s emphasis on PME, operational assignments, and self-development raise our leadership lid?',
        'What is one lid that commonly limits leaders in our unit right now?',
        'What conditions in a unit help leaders grow, and what conditions keep them stuck?'
    ],
    [
        'Leadership development is a duty, not a hobby.',
        'The first step is honest self-assessment, followed by one clear action.'
    ]
)

concept_slide(
    'Law of Influence | Expanded',
    'Original framework idea: The true measure of leadership is influence, nothing more, nothing less.',
    'Doctrine tie-in and explanation', [
        'ADP 6-22 defines leadership as influencing people by providing purpose, direction, and motivation.',
        'This means influence is already central to Army leadership, not just a business-book idea.',
        'Rank can compel compliance, but influence creates commitment, initiative, and follow-through.',
        'The best leaders explain intent, enforce standards fairly, and make Soldiers feel led, not managed.'
    ],
    'What to emphasize in the room', [
        'Influence grows when leaders are competent, consistent, and personally invested in their people.',
        'It weakens when leaders hide behind rank or only show up to correct problems.',
        'Soldiers can tell the difference between authority and leadership very quickly.'
    ],
    'A real leader knows the difference between being a boss and being a leader.',
    'John C. Maxwell'
)

discussion_slide(
    'Law of Influence | Facilitated Discussion',
    'This is usually the easiest concept for Soldiers to connect with because everyone has seen both good and bad examples.',
    [
        'Start with a simple distinction: authority may get movement, but influence gets effort and ownership.',
        'Invite examples of leaders who built influence without raising their voice or pulling rank.',
        'Tie responses back to purpose, direction, and motivation.'
    ],
    [
        'How do we build influence beyond our rank?',
        'What specific leader behaviors make Soldiers want to follow instead of just comply?',
        'How does poor communication damage influence even when authority is still intact?'
    ],
    [
        'Influence is earned daily through competence, clarity, fairness, and example.',
        'A leader who cannot influence will eventually lean too hard on position.'
    ]
)

concept_slide(
    'Law of Solid Ground | Expanded',
    'Original framework idea: Trust is the foundation of leadership.',
    'Doctrine tie-in and explanation', [
        'AR 600-100 and DA PAM 165-19 reinforce that trust is foundational to the profession and to relationships inside the formation.',
        'Trust depends on character, competence, and commitment, not just intentions.',
        'When trust is high, communication is more honest, standards are more credible, and teams are more resilient.',
        'When trust is low, even correct orders are received with doubt or disengagement.'
    ],
    'What to emphasize in the room', [
        'Trust is built in small moments: fairness, follow-through, truthfulness, and moral courage.',
        'It is usually eroded by inconsistency, favoritism, weak accountability, or broken promises.',
        'Rebuilding trust takes time and visible behavior change.'
    ],
    'Trust is the bedrock upon which we ground our relationship with the American people.',
    'DA PAM 165-19'
)

discussion_slide(
    'Law of Solid Ground | Facilitated Discussion',
    'This is the most important section if you want the room to connect leadership theory to unit climate.',
    [
        'Ask for examples of trust-builders and trust-breakers in real formations.',
        'Keep the discussion professional and avoid turning it into a complaint session.',
        'If needed, redirect from personalities to patterns of behavior.'
    ],
    [
        'What actions erode trust fastest in a unit?',
        'How can leaders be more intentional about building and maintaining trust with Soldiers, peers, and superiors?',
        'If trust is damaged, what does rebuilding it actually require?'
    ],
    [
        'Trust is not abstract. Soldiers feel it in fairness, honesty, competence, and accountability.',
        'Without trust, leadership influence collapses over time.'
    ]
)

exercise_slide()
icebreakers_slide()
close_slide()
refs_slide()

prs.save(OUT_PPTX)

notes = f'''# Session 1 LPD Framework Expanded - Speaker Notes

## Deliverables
- Deck: `{OUT_PPTX}`
- Notes: `{OUT_NOTES}`

## Intent
This revised version stays closer to the original spreadsheet framework. Instead of turning the content into a generic presentation, it expands the framework into a facilitator-ready LPD with more talk tracks and clearer transitions.

## Suggested flow
1. **Title**
   - Add your unit, rank, date, and any local branding.
   - Opening line: "This session is about the leadership foundations that determine whether Soldiers merely comply or actually trust and follow."

2. **How to use this expanded framework**
   - Tell the room this is a discussion-based LPD.
   - Set the tone: practical, honest, and applicable.

3. **Original framework expanded**
   - Use this to show the lesson still comes directly from the framework.
   - Key line: "Maxwell gives us language. Army doctrine gives us the standard."

4. **Law of the Lid slides**
   - Focus on self-development and how leader limitations affect the unit.
   - Good local question: "Where do we most often see junior leaders hit a ceiling in our formation?"

5. **Law of Influence slides**
   - Emphasize that influence is how leadership actually works in the Army.
   - Good local question: "What made you trust or follow a leader beyond just their rank?"

6. **Law of Solid Ground slides**
   - This is where the room usually gets most engaged.
   - Good local question: "What behaviors cause Soldiers to stop believing what leaders say?"

7. **Exercises**
   - Use the self-assessment if you want reflection and commitment.
   - Use the role-play if you want energy and visible participation.

8. **Icebreakers**
   - Use only one.
   - Best fast option: One Word for Leadership.

9. **Closeout**
   - End with visible, near-term action.
   - If possible, revisit their commitments later.

## Minimal edits before presenting
- Add your unit and name on the title slide.
- Drop in 1 to 3 examples from your own formation.
- Adjust wording for audience maturity and rank.
- If you need a shorter version, cut one discussion slide per concept and keep only one exercise.

## What changed from the first version
- The new deck is tied more tightly to the spreadsheet framework.
- Each original concept now has expanded explanation plus a facilitator discussion slide.
- The overall tone is more like an LPD framework you can brief from, not just a polished summary deck.
'''
OUT_NOTES.write_text(notes)
print(OUT_PPTX)
print(OUT_NOTES)
