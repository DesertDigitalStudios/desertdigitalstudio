from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT_DIR = Path('deliverables')
OUT_DIR.mkdir(exist_ok=True)
OUT_PPTX = OUT_DIR / 'Session_1_Foundation_of_Leadership_Keynote_Ready.pptx'
OUT_NOTES = OUT_DIR / 'Session_1_Foundation_of_Leadership_Speaker_Notes.md'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette
BG = RGBColor(245, 246, 242)
DARK = RGBColor(29, 40, 34)
OLIVE = RGBColor(77, 100, 72)
OLIVE_LIGHT = RGBColor(202, 213, 191)
TAN = RGBColor(230, 223, 203)
ACCENT = RGBColor(135, 111, 74)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(145, 57, 49)

TITLE_FONT = 'Arial'
BODY_FONT = 'Arial'


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title, subtitle=None, band_color=DARK):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = band_color
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(9.6), Inches(0.42))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = TITLE_FONT
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.level = 0
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = BODY_FONT
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(220, 226, 217)


def add_footer(slide, text='LPD Session 1 | Leadership Foundations'):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = OLIVE
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.02), Inches(6), Inches(0.22))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(90, 95, 88)


def add_bullets(slide, left, top, width, height, bullets, font_size=20, color=DARK, bullet_color=None, spacing=1.1):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(2)
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        p.line_spacing = spacing
        p.bullet = True if level >= 0 else False
        r = p.add_run()
        r.text = text
        r.font.name = BODY_FONT
        r.font.size = Pt(font_size - (2 * level))
        r.font.color.rgb = color
        if bullet_color:
            p.font.color.rgb = bullet_color
    return box


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_callout(slide, left, top, width, height, header, bullets, fill_color=WHITE, header_color=OLIVE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = OLIVE_LIGHT
    header_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.45))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = header_color
    header_box.line.fill.background()
    hb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.08), width - Inches(0.3), Inches(0.22))
    tf = hb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = header
    r.font.name = TITLE_FONT
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = WHITE
    add_bullets(slide, left + Inches(0.18), top + Inches(0.56), width - Inches(0.32), height - Inches(0.66), bullets, font_size=16)


def add_title_slide(title, subtitle, details):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, color=DARK)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.8), Inches(0.18), Inches(4.9))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = OLIVE_LIGHT
    stripe.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(10.6), Inches(2.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = TITLE_FONT
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = BODY_FONT
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(220, 226, 217)
    detail_box = slide.shapes.add_textbox(Inches(1.05), Inches(4.9), Inches(6.6), Inches(1.2))
    dtf = detail_box.text_frame
    for i, line in enumerate(details):
        p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = BODY_FONT
        r.font.size = Pt(16 if i == 0 else 14)
        r.font.color.rgb = RGBColor(228, 232, 225)
        if i == 0:
            r.font.bold = True
    quote_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(4.45), Inches(4.25), Inches(1.55))
    quote_shape.fill.solid()
    quote_shape.fill.fore_color.rgb = RGBColor(43, 57, 50)
    quote_shape.line.color.rgb = OLIVE
    q = slide.shapes.add_textbox(Inches(8.45), Inches(4.72), Inches(3.8), Inches(1.0))
    qtf = q.text_frame
    qtf.word_wrap = True
    p = qtf.paragraphs[0]
    r = p.add_run()
    r.text = 'Leadership is influence, grounded in trust.'
    r.font.name = BODY_FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    add_textbox(slide, Inches(8.45), Inches(5.46), Inches(3.4), Inches(0.28), 'Maxwell + Army Doctrine Crosswalk', font_size=11, color=RGBColor(205, 214, 203))


def add_content_slide(title, subtitle, left_title, left_bullets, right_header, right_bullets, footer='LPD Session 1 | Leadership Foundations'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, title, subtitle)
    section = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.3), Inches(7.2), Inches(5.35))
    section.fill.solid()
    section.fill.fore_color.rgb = WHITE
    section.line.color.rgb = OLIVE_LIGHT
    add_textbox(slide, Inches(0.82), Inches(1.58), Inches(6.5), Inches(0.35), left_title, font_size=18, bold=True, color=OLIVE)
    add_bullets(slide, Inches(0.82), Inches(2.0), Inches(6.45), Inches(4.35), left_bullets, font_size=18)
    add_callout(slide, Inches(7.95), Inches(1.45), Inches(4.75), Inches(5.15), right_header, right_bullets, fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
    add_footer(slide, footer)


def add_agenda_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, 'Purpose, Outcomes, and Flow', 'Built as a ready-to-edit LPD brief for Soldiers and junior leaders')
    add_callout(slide, Inches(0.65), Inches(1.5), Inches(5.9), Inches(4.6), 'Learning objectives', [
        'Connect Maxwell\'s Laws of the Lid, Influence, and Solid Ground to Army doctrine.',
        'Identify behaviors that raise or lower a leader\'s effectiveness.',
        'Discuss practical ways to build influence beyond rank.',
        'Reinforce trust as the foundation of cohesive teams and mission success.',
        'Leave with one 30-day action to improve leadership.'
    ], fill_color=WHITE, header_color=OLIVE)
    add_callout(slide, Inches(6.8), Inches(1.5), Inches(5.85), Inches(4.6), 'Recommended flow (45 to 60 min)', [
        '5 to 10 min: Icebreaker and framing question',
        '10 min: Why leadership foundations matter in the Army',
        '15 to 20 min: Maxwell and doctrine crosswalk',
        '10 to 15 min: Discussion and practical exercise',
        '5 min: Commitments, takeaways, and close'
    ], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
    add_textbox(slide, Inches(0.75), Inches(6.35), Inches(11.9), Inches(0.28), 'Tip: personalize this deck with your unit, recent examples, and one local leadership success story.', font_size=12, color=RGBColor(95, 100, 93), bold=False)
    add_footer(slide)


def add_crosswalk_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, 'Maxwell to Army Doctrine Crosswalk', 'The lesson works because the concepts match what Army doctrine already expects from leaders.')
    table = slide.shapes.add_table(4, 3, Inches(0.75), Inches(1.65), Inches(11.9), Inches(3.35)).table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.7)
    table.columns[2].width = Inches(5.2)
    headers = ['Maxwell principle', 'Primary Army references', 'Army leadership concept']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = OLIVE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = TITLE_FONT
                r.font.size = Pt(15)
                r.font.bold = True
                r.font.color.rgb = WHITE
    rows = [
        ('Law of the Lid', 'ADP 6-22, FM 6-22', 'Leader development and self-development raise a leader\'s ceiling and improve unit effectiveness.'),
        ('Law of Influence', 'ADP 6-22, FM 1', 'Leadership is the process of influencing people by providing purpose, direction, and motivation.'),
        ('Law of Solid Ground', 'AR 600-100, DA PAM 165-19', 'Trust is the bedrock of the profession and depends on character, competence, and commitment.')
    ]
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else RGBColor(248, 248, 244)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = BODY_FONT
                    run.font.size = Pt(14)
                    run.font.color.rgb = DARK
    add_callout(slide, Inches(0.95), Inches(5.35), Inches(11.3), Inches(1.05), 'Use this framing line', ['These are not outside ideas replacing Army leadership. They are a practical way to explain the same standards we already expect in our formation.'], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
    add_footer(slide)


def add_two_column_text_slide(title, subtitle, left_header, left_bullets, right_header, right_bullets, footer='LPD Session 1 | Leadership Foundations'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, title, subtitle)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3))
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = OLIVE_LIGHT
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(251, 250, 247)
    right.line.color.rgb = OLIVE_LIGHT
    add_textbox(slide, Inches(0.95), Inches(1.78), Inches(4.9), Inches(0.35), left_header, font_size=18, bold=True, color=OLIVE)
    add_bullets(slide, Inches(0.95), Inches(2.15), Inches(5.15), Inches(4.35), left_bullets, font_size=18)
    add_textbox(slide, Inches(7.05), Inches(1.78), Inches(4.9), Inches(0.35), right_header, font_size=18, bold=True, color=ACCENT)
    add_bullets(slide, Inches(7.05), Inches(2.15), Inches(5.1), Inches(4.35), right_bullets, font_size=17)
    add_footer(slide, footer)


def add_quote_banner(slide, quote, attribution):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.8), Inches(0.82))
    box.fill.solid()
    box.fill.fore_color.rgb = TAN
    box.line.color.rgb = RGBColor(211, 198, 168)
    tb = slide.shapes.add_textbox(Inches(1.02), Inches(6.0), Inches(11.35), Inches(0.42))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f'"{quote}"'
    r.font.name = BODY_FONT
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = attribution
    r2.font.name = BODY_FONT
    r2.font.size = Pt(11)
    r2.font.color.rgb = RGBColor(92, 83, 70)


def add_exercise_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, 'Practical Application Exercises', 'Use one exercise if time is short, or both if you want stronger engagement and ownership.')
    add_callout(slide, Inches(0.7), Inches(1.55), Inches(5.85), Inches(4.95), 'Exercise 1: The Lid self-assessment', [
        'Have leaders rate themselves against ADP 6-22 competencies: Leads, Develops, Achieves.',
        'Ask: What is the single biggest lid holding back my effectiveness right now?',
        'Ask: Which Army publication or mentor can help me improve?',
        'Ask: What one action will I take in the next 30 days?',
        'Close with small-group sharing of goals, not private weaknesses.'
    ], fill_color=WHITE, header_color=OLIVE)
    add_callout(slide, Inches(6.8), Inches(1.55), Inches(5.85), Inches(4.95), 'Exercise 2: Influence beyond rank', [
        'Pair up: one squad leader, one skeptical new Soldier.',
        'Scenario: a last-minute clean-up detail cancels early release.',
        'Task: earn buy-in without saying, "because I said so."',
        'After the role-play, ask what words or behaviors actually influenced followership.',
        'AAR focus: clarity, empathy, standards, and leader presence.'
    ], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
    add_footer(slide)


def add_icebreaker_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, 'Icebreaker Options', 'Pick one based on time, audience maturity, and how much participation you want up front.')
    add_two_column_text_slide_content = False
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(5.9), Inches(5.25))
    left.fill.solid(); left.fill.fore_color.rgb = WHITE; left.line.color.rgb = OLIVE_LIGHT
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.5), Inches(5.85), Inches(5.25))
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor(251, 250, 247); right.line.color.rgb = OLIVE_LIGHT
    add_textbox(slide, Inches(0.95), Inches(1.8), Inches(5), Inches(0.3), 'Best quick options', font_size=18, bold=True, color=OLIVE)
    add_bullets(slide, Inches(0.95), Inches(2.15), Inches(5.1), Inches(4.15), [
        'One Word for Leadership, 3 to 5 min. Ask each person for one word and one sentence why.',
        'Leadership If-Then Statement, 3 to 5 min. Example: If I prioritize Soldiers\' development, then...',
        'Two Truths and a Lie, leadership edition, 5 to 7 min for a lighter start.'
    ], font_size=17)
    add_textbox(slide, Inches(7.0), Inches(1.8), Inches(5), Inches(0.3), 'Best deeper options', font_size=18, bold=True, color=ACCENT)
    add_bullets(slide, Inches(7.0), Inches(2.15), Inches(5.1), Inches(4.15), [
        'Leadership Superpower, 7 to 10 min. Good for surfacing values and humor.',
        'Most Valuable Leadership Lesson Learned, 10 to 15 min. Best for experienced groups.',
        'Recommendation: use One Word for Leadership if you need fast participation and a clean transition into trust and influence.'
    ], font_size=17)
    add_footer(slide)


def add_takeaways_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, 'Close Strong', 'End with action, not just discussion.')
    add_callout(slide, Inches(0.75), Inches(1.6), Inches(5.95), Inches(4.9), 'Three takeaways', [
        'A leader\'s lid can raise or lower the performance of the whole team.',
        'Influence matters more than title when you need commitment, discipline, and initiative.',
        'Trust is fragile, and every daily interaction either builds it or erodes it.'
    ], fill_color=WHITE, header_color=OLIVE)
    add_callout(slide, Inches(6.9), Inches(1.6), Inches(5.7), Inches(4.9), '30-day leader challenge', [
        'Identify one leadership lid you will raise this month.',
        'Use explanation and example before authority at least once this week.',
        'Do one visible action that builds trust with your Soldiers.',
        'Write the action down and follow up at the next LPD or counseling session.'
    ], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
    add_textbox(slide, Inches(0.85), Inches(6.45), Inches(11.7), Inches(0.25), 'Suggested closing question: What will your Soldiers see from you this week that proves this lesson mattered?', font_size=13, color=RGBColor(96, 100, 93), bold=True)
    add_footer(slide)


def add_references_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, 'References', 'Primary sources used to build the lesson')
    refs = [
        'John C. Maxwell, Leadership 101, especially the Laws of the Lid, Influence, and Solid Ground.',
        'ADP 6-22, Army Leadership and the Profession.',
        'FM 6-22, Developing Leaders.',
        'FM 1, The Army.',
        'AR 600-100, Army Profession and Leadership Policy.',
        'DA PAM 165-19, Moral Leadership.',
        'Supporting references from the source spreadsheet: ADP 1, DA PAM 600-25, TC 7-22.7, and ADP 7-0.'
    ]
    add_bullets(slide, Inches(0.9), Inches(1.7), Inches(11.4), Inches(4.9), refs, font_size=18)
    add_callout(slide, Inches(0.9), Inches(5.75), Inches(11.35), Inches(0.9), 'Final edit reminder', ['Swap in unit examples, add your rank and unit on slide 1, and tailor the discussion prompts to your Soldiers.'], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
    add_footer(slide)


# Build slides
add_title_slide(
    'Session 1: The Foundation of Leadership',
    'Influence and Trust | Leader Professional Development (LPD)',
    [
        'Keynote-ready deck built from the provided lesson spreadsheet',
        'Crosswalks John C. Maxwell Leadership 101 with Army doctrine and leader development guidance',
        'Prepared for easy final editing before delivery'
    ]
)
add_agenda_slide()

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, 'Why This Lesson Matters', 'Army leadership is not just authority. It is the consistent ability to influence people and earn trust.')
add_content_slide_content = False
main = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.45), Inches(7.0), Inches(4.2))
main.fill.solid(); main.fill.fore_color.rgb = WHITE; main.line.color.rgb = OLIVE_LIGHT
add_textbox(slide, Inches(0.95), Inches(1.8), Inches(6.3), Inches(0.3), 'Why it matters in a formation', font_size=18, bold=True, color=OLIVE)
add_bullets(slide, Inches(0.95), Inches(2.15), Inches(6.25), Inches(3.1), [
    'Mission success depends on purpose, direction, motivation, and disciplined initiative.',
    'Soldiers may obey rank, but they commit to leaders they trust and respect.',
    'Every leader sets conditions for growth, standards, and climate.',
    'Poor leadership does not stay personal. It lowers the ceiling for the whole team.'
], font_size=18)
add_callout(slide, Inches(8.0), Inches(1.65), Inches(4.55), Inches(2.0), 'Anchor idea', [
    'A leader who develops self, influences well, and protects trust gives the unit an operational advantage.'
], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
add_callout(slide, Inches(8.0), Inches(3.95), Inches(4.55), Inches(1.45), 'Use this question to open discussion', [
    'When Soldiers talk about good leaders they have had, what do they usually remember most?' ], fill_color=WHITE, header_color=OLIVE)
add_footer(slide)
add_quote_banner(slide, 'A real leader knows the difference between being a boss and being a leader.', 'John C. Maxwell')

add_crosswalk_slide()

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, 'Law of the Lid', 'Leadership ability determines a person\'s level of effectiveness.')
left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(6.1), Inches(5.05))
left.fill.solid(); left.fill.fore_color.rgb = WHITE; left.line.color.rgb = OLIVE_LIGHT
add_textbox(slide, Inches(0.95), Inches(1.8), Inches(5.2), Inches(0.3), 'Army connection', font_size=18, bold=True, color=OLIVE)
add_bullets(slide, Inches(0.95), Inches(2.15), Inches(5.35), Inches(3.9), [
    'ADP 6-22 and FM 6-22 emphasize continuous leader development and self-development.',
    'Leaders are responsible for raising their own ceiling and helping subordinates raise theirs.',
    'PME, repetition, coaching, counseling, and honest feedback all raise the lid.'
], font_size=18)
add_callout(slide, Inches(7.1), Inches(1.55), Inches(5.45), Inches(3.0), 'Discussion prompts', [
    'How does the Army\'s emphasis on PME, assignments, and self-development raise a leader\'s lid?',
    'What is one lid that commonly limits junior leaders in our formation?',
    'How can we build a climate where growth is expected instead of optional?'
], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
add_callout(slide, Inches(7.1), Inches(4.85), Inches(5.45), Inches(1.65), 'Suggested leader action', [
    'Have each person identify one competency to improve in the next 30 days.'
], fill_color=WHITE, header_color=OLIVE)
add_footer(slide)
add_quote_banner(slide, 'The lower an individual\'s ability to lead, the lower the lid on his potential. The higher the leadership, the greater the effectiveness.', 'John C. Maxwell')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, 'Law of Influence', 'The true measure of leadership is influence, nothing more, nothing less.')
left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(6.05), Inches(5.05))
left.fill.solid(); left.fill.fore_color.rgb = WHITE; left.line.color.rgb = OLIVE_LIGHT
add_textbox(slide, Inches(0.95), Inches(1.8), Inches(5.2), Inches(0.3), 'Army connection', font_size=18, bold=True, color=OLIVE)
add_bullets(slide, Inches(0.95), Inches(2.15), Inches(5.35), Inches(3.9), [
    'ADP 6-22 defines leadership as influencing people by providing purpose, direction, and motivation.',
    'Rank gives authority. Influence earns trust, effort, initiative, and follow-through.',
    'Leaders build influence when they explain the why, set the example, and show they know their people.'
], font_size=18)
add_callout(slide, Inches(7.05), Inches(1.55), Inches(5.5), Inches(2.65), 'Practical ways to build influence', [
    'Be consistent and keep standards fair.',
    'Communicate intent, not just orders.',
    'Know your Soldiers and their motivations.',
    'Solve problems instead of hiding behind rank.'
], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
add_callout(slide, Inches(7.05), Inches(4.55), Inches(5.5), Inches(1.95), 'Discussion prompt', [
    'How do we build influence beyond our rank so Soldiers want to follow, not just comply?' ], fill_color=WHITE, header_color=OLIVE)
add_footer(slide)
add_quote_banner(slide, 'A real leader knows the difference between being a boss and being a leader.', 'John C. Maxwell')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_top_band(slide, 'Law of Solid Ground', 'Trust is the foundation of leadership.')
left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(6.1), Inches(5.05))
left.fill.solid(); left.fill.fore_color.rgb = WHITE; left.line.color.rgb = OLIVE_LIGHT
add_textbox(slide, Inches(0.95), Inches(1.8), Inches(5.2), Inches(0.3), 'Army connection', font_size=18, bold=True, color=OLIVE)
add_bullets(slide, Inches(0.95), Inches(2.15), Inches(5.35), Inches(3.9), [
    'AR 600-100 and DA PAM 165-19 frame trust as foundational to the Army profession.',
    'Trust grows from character, competence, and commitment, not from slogans.',
    'Without trust, standards slip, communication narrows, and cohesion breaks down.'
], font_size=18)
add_callout(slide, Inches(7.05), Inches(1.55), Inches(5.5), Inches(2.7), 'Trust builders and trust breakers', [
    'Build trust: honesty, competence, follow-through, fairness, accountability.',
    'Break trust: favoritism, inconsistency, gossip, weak moral courage, empty promises.'
], fill_color=RGBColor(251, 250, 247), header_color=ACCENT)
add_callout(slide, Inches(7.05), Inches(4.55), Inches(5.5), Inches(1.95), 'Discussion prompt', [
    'What actions erode trust fastest in a unit, and how do leaders rebuild it after failure?' ], fill_color=WHITE, header_color=OLIVE)
add_footer(slide)
add_quote_banner(slide, 'Trust is the bedrock upon which we ground our relationship with the American people.', 'DA PAM 165-19')

add_exercise_slide()
add_icebreaker_slide()
add_takeaways_slide()
add_references_slide()

prs.save(OUT_PPTX)

notes = f"""# Session 1: The Foundation of Leadership - Speaker Notes

## Files
- Deck: `{OUT_PPTX}`
- Purpose: ready-to-edit LPD deck for Soldiers and junior leaders

## Suggested delivery flow
1. **Title slide**
   - Add your unit, rank, and date before presenting.
   - Open with: *"Today is about the foundations that make people actually want to follow a leader: growth, influence, and trust."*

2. **Purpose, outcomes, and flow**
   - Set expectations: this is discussion-based, not a one-way lecture.
   - If time is tight, tell them up front you will use one icebreaker and one practical exercise.

3. **Why this lesson matters**
   - Connect the topic to mission success, climate, and standards.
   - Use one recent example from your unit if possible.

4. **Crosswalk slide**
   - Frame Maxwell as a practical language aid, not a replacement for Army doctrine.
   - Key line: *"These principles reinforce what the Army already expects from leaders."*

5. **Law of the Lid**
   - Push reflection without making it feel like a character attack.
   - Strong transition: *"If our leadership ceiling is low, our unit feels it fast."*

6. **Law of Influence**
   - Make the point that rank can force movement, but influence creates commitment.
   - Ask for examples of leaders who influenced them without yelling or threatening.

7. **Law of Solid Ground**
   - Keep this practical. Ask what breaks trust in real units.
   - Good follow-up question: *"What do Soldiers notice first when trust starts slipping?"*

8. **Exercises**
   - Pick the self-assessment if you want reflection.
   - Pick the role-play if you want energy and discussion.
   - If time allows, use both and end with a short AAR.

9. **Icebreakers**
   - Best fast option: *One Word for Leadership*.
   - Best deeper option: *Most Valuable Leadership Lesson Learned*.

10. **Takeaways / close**
   - End by making everyone commit to one visible action.
   - If appropriate, collect their 30-day action and revisit it later.

## Minimal edits I recommend before you present
- Put your unit name and date on slide 1.
- Add one example from your formation on slides 5 to 8.
- Tighten any language to match your audience, especially if they are younger Soldiers.
- If your unit prefers shorter decks, you can cut the icebreaker slide and keep the notes for yourself.

## Source summary from the spreadsheet
- Session title: **The Foundation of Leadership - Influence and Trust**
- Main concepts: **Law of the Lid, Law of Influence, Law of Solid Ground**
- Army references: **ADP 6-22, FM 6-22, FM 1, AR 600-100, DA PAM 165-19**
- Included activities: **self-assessment, role-play, multiple icebreakers**
"""
OUT_NOTES.write_text(notes)
print(f'Created: {OUT_PPTX}')
print(f'Created: {OUT_NOTES}')
